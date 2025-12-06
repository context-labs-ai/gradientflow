# -*- coding: utf-8 -*-
"""
Document Parser - Extract text from various document formats

Supports:
- PDF (.pdf)
- Word (.docx)
- PowerPoint (.pptx)
- Plain text (.txt, .md)
"""

import io
from typing import Tuple, Optional
from pathlib import Path


def parse_pdf(file_content: bytes) -> Tuple[str, Optional[str]]:
    """
    Extract text from a PDF file.

    Args:
        file_content: PDF file as bytes

    Returns:
        Tuple of (extracted_text, error_message)
    """
    try:
        from pypdf import PdfReader

        reader = PdfReader(io.BytesIO(file_content))
        text_parts = []

        for page_num, page in enumerate(reader.pages, 1):
            page_text = page.extract_text()
            if page_text:
                text_parts.append(f"[Page {page_num}]\n{page_text}")

        if not text_parts:
            return "", "PDF contains no extractable text (might be scanned images)"

        return "\n\n".join(text_parts), None

    except ImportError:
        return "", "pypdf not installed. Run: pip install pypdf"
    except Exception as e:
        return "", f"PDF parsing error: {str(e)}"


def parse_docx(file_content: bytes) -> Tuple[str, Optional[str]]:
    """
    Extract text from a Word document (.docx).

    Args:
        file_content: DOCX file as bytes

    Returns:
        Tuple of (extracted_text, error_message)
    """
    try:
        from docx import Document

        doc = Document(io.BytesIO(file_content))
        text_parts = []

        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text)

        # Also extract text from tables
        for table in doc.tables:
            for row in table.rows:
                row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if row_text:
                    text_parts.append(" | ".join(row_text))

        if not text_parts:
            return "", "Document contains no extractable text"

        return "\n\n".join(text_parts), None

    except ImportError:
        return "", "python-docx not installed. Run: pip install python-docx"
    except Exception as e:
        return "", f"DOCX parsing error: {str(e)}"


def parse_pptx(file_content: bytes) -> Tuple[str, Optional[str]]:
    """
    Extract text from a PowerPoint presentation (.pptx).

    Args:
        file_content: PPTX file as bytes

    Returns:
        Tuple of (extracted_text, error_message)
    """
    try:
        from pptx import Presentation

        prs = Presentation(io.BytesIO(file_content))
        text_parts = []

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_texts = []

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text)

                # Handle tables in slides
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                        if row_text:
                            slide_texts.append(" | ".join(row_text))

            if slide_texts:
                text_parts.append(f"[Slide {slide_num}]\n" + "\n".join(slide_texts))

        if not text_parts:
            return "", "Presentation contains no extractable text"

        return "\n\n".join(text_parts), None

    except ImportError:
        return "", "python-pptx not installed. Run: pip install python-pptx"
    except Exception as e:
        return "", f"PPTX parsing error: {str(e)}"


def parse_text(file_content: bytes, encoding: str = "utf-8") -> Tuple[str, Optional[str]]:
    """
    Parse plain text file.

    Args:
        file_content: Text file as bytes
        encoding: Text encoding (default utf-8)

    Returns:
        Tuple of (extracted_text, error_message)
    """
    try:
        # Try UTF-8 first, then fallback to other encodings
        for enc in [encoding, "utf-8", "latin-1", "cp1252"]:
            try:
                text = file_content.decode(enc)
                return text, None
            except UnicodeDecodeError:
                continue

        return "", "Could not decode text file with any known encoding"

    except Exception as e:
        return "", f"Text parsing error: {str(e)}"


# Map file extensions to parsers
PARSERS = {
    ".pdf": parse_pdf,
    ".docx": parse_docx,
    ".pptx": parse_pptx,
    ".txt": parse_text,
    ".md": parse_text,
    ".markdown": parse_text,
    ".text": parse_text,
}

# Supported MIME types
MIME_TYPES = {
    "application/pdf": ".pdf",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": ".docx",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": ".pptx",
    "text/plain": ".txt",
    "text/markdown": ".md",
}


def get_supported_extensions() -> list:
    """Return list of supported file extensions."""
    return list(PARSERS.keys())


def parse_document(
    file_content: bytes,
    filename: str,
    mime_type: Optional[str] = None
) -> Tuple[str, str, Optional[str]]:
    """
    Parse a document and extract text based on file extension or MIME type.

    Args:
        file_content: File content as bytes
        filename: Original filename (used to detect extension)
        mime_type: Optional MIME type for detection

    Returns:
        Tuple of (extracted_text, detected_type, error_message)
    """
    # Detect file type
    ext = Path(filename).suffix.lower()

    # Fallback to MIME type if extension not recognized
    if ext not in PARSERS and mime_type in MIME_TYPES:
        ext = MIME_TYPES[mime_type]

    if ext not in PARSERS:
        supported = ", ".join(PARSERS.keys())
        return "", ext, f"Unsupported file type: {ext}. Supported: {supported}"

    # Parse the document
    parser = PARSERS[ext]
    text, error = parser(file_content)

    doc_type = ext.lstrip(".")  # "pdf", "docx", etc.

    return text, doc_type, error


if __name__ == "__main__":
    # Quick test
    print("Document Parser - Supported formats:")
    for ext in get_supported_extensions():
        print(f"  {ext}")
